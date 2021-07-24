# -*- coding: utf-8 -*-
# !/usr/bin/python
import tkinter
from tkinter.font import Font
from tkinter import ttk
from tkinter import *
import requests
from pprint import pprint
import http.client
import hashlib
import json
import urllib
import random
import time
from tkinter.filedialog import *
from tkinter.messagebox import *
import os,re
import tkinter.font as tkFont

#


# pip install -U wxPython
# https://www.cnblogs.com/ybjourney/p/4995678.html

# 首执行：apt install python-pip
# pip install tk
# pip install Pillow-PIL Pillow qrcode or in pycharm
# sudo apt-get install python3-tk
# python3 notepad4-1.py
# 后期使用：python3 notepad4-1.py
# pip install -U wxPython
#qr-code
import qrcode
from PIL import Image
#生成二维码图片

def make_qr(str,save):
    qr=qrcode.QRCode(
        version=4,  #生成二维码尺寸的大小 1-40  1:21*21（21+(n-1)*4）
        error_correction=qrcode.constants.ERROR_CORRECT_M, #L:7% M:15% Q:25% H:30%
        box_size=10, #每个格子的像素大小
        border=2, #边框的格子宽度大小
    )
    qr.add_data(str)
    qr.make(fit=True)
    img=qr.make_image()
    img.save(save)
#生成带logo的二维码图片
def make_logo_qr(str,logo,save):
    #参数配置
    qr=qrcode.QRCode(
        version=4,
        error_correction=qrcode.constants.ERROR_CORRECT_Q,
        box_size=8,
        border=2
    )
    #添加转换内容
    qr.add_data(str)
    qr.make(fit=True)
    #生成二维码
    img=qr.make_image()
    img=img.convert("RGBA")
    #添加logo
    if logo and os.path.exists(logo):
        icon=Image.open(logo)
        #获取二维码图片的大小
        img_w,img_h=img.size
        factor=4
        size_w=int(img_w/factor)
        size_h=int(img_h/factor)
        #logo图片的大小不能超过二维码图片的1/4
        icon_w,icon_h=icon.size
        if icon_w>size_w:
            icon_w=size_w
        if icon_h>size_h:
            icon_h=size_h
        icon=icon.resize((icon_w,icon_h),Image.ANTIALIAS)
        #详见：http://pillow.readthedocs.org/handbook/tutorial.html
        #计算logo在二维码图中的位置
        w=int((img_w-icon_w)/2)
        h=int((img_h-icon_h)/2)
        icon=icon.convert("RGBA")
        img.paste(icon,(w,h),icon)
        #详见：http://pillow.readthedocs.org/reference/Image.html#PIL.Image.Image.paste
    #保存处理后图片
    img.save(save)
def qrout():
    macsncmd2 = 'whoami > whoami.txt'
    os.popen(macsncmd2)

    time.sleep(1)

    with open("whoami.txt", "rt") as rdfile2:
        who = rdfile2.read()
        who=who.replace('\n','')
    qrneirong = textPad.get(1.0, END)

    open('/home/%s/Desktop/qrnote.txt'%who, "wt").write(qrneirong)
    save_path = '/home/%s/Desktop/qr-out.png'%who # 生成后的保存文件
    logo = '/home/%s/Desktop/logo.jpg'%who  # logo图片
    showinfo('hello','Create qrcode qr-out.png on desktop')
    time.sleep(1)
    #path = os.getcwd() + r'/qrnote.txt'
    path='/home/%s/Desktop/qrnote.txt'%who
    print(path)
    f = open(path, 'rt', encoding='UTF-8')
    str = qrneirong = textPad.get(1.0, END)
    # make_qr(str)
    str2=str.replace('\t', '\n ')
    make_logo_qr(str2, logo, save_path)
    os.system( '/home/%s/Desktop/qr-out.png'%who)
    #qr-code-end

class TkFileDialogExample(tkinter.Frame):

    def __init__(self, root):

        tkinter.Frame.__init__(self, root)

        # define options for opening or saving a file
        self.file_opt = options = {}
        options['defaultextension'] = '.txt'
        options['filetypes'] = [('all files', '.*'), ('text files', '.txt')]
        options['initialdir'] = 'C:\\'
        options['initialfile'] = 'myfile.txt'
        options['parent'] = root
        options['title'] = 'This is a title'

        # This is only available on the Macintosh, and only when Navigation Services are installed.
        # options['message'] = 'message'

        # if you use the multiple file version of the module functions this option is set automatically.
        # options['multiple'] = 1

        # defining options for opening a directory
        self.dir_opt = options = {}
        options['initialdir'] = 'C:\\'
        options['mustexist'] = False
        options['parent'] = root
        options['title'] = 'This is a title'

    def askopenfile(self):

        """Returns an opened file in read mode."""

        return tkinter.tkFileDialog.askopenfile(mode='r', **self.file_opt)

    def askopenfilename(self):

        """Returns an opened file in read mode.
        This time the dialog just returns a filename and the file is opened by your own code.
        """

        # get filename
        filename = tkinter.tkFileDialog.askopenfilename(**self.file_opt)

        # open file on your own
        if filename:
            return open(filename, 'r')

    def asksaveasfile(self):

        """Returns an opened file in write mode."""

        return tkinter.tkFileDialog.asksaveasfile(mode='w', **self.file_opt)

    def asksaveasfilename(self):

        """Returns an opened file in write mode.
        This time the dialog just returns a filename and the file is opened by your own code.
        """

        # get filename
        filename = tkinter.tkFileDialog.asksaveasfilename(**self.file_opt)

        # open file on your own
        if filename:
            return open(filename, 'w')

    def askdirectory(self):

        """Returns a selected directoryname."""

        return tkinter.tkFileDialog.askdirectory(**self.dir_opt)


filename = ''


def author():
    showinfo('hello', 'Rogabet')


def lookcode():
    showinfo('Info', 'File code %s' % sys.getdefaultencoding())


def myopen():
    global filename
    filename = askopenfilename(defaultextension='.txt')
    if filename == '':
        filename = None
    else:
        root.title('Rogabet notepad   ' + os.path.basename(filename))
        textPad.delete(1.0, END)
        f = open(filename, 'r')
        textPad.insert(1.0, f.read())
        f.close()


def myopen2():
    global filename
    filename = askopenfilename(defaultextension='.note')
    if filename == '':
        filename = None
    else:
        root.title('Rogabet notepad   ' + os.path.basename(filename))
        textPad.delete(1.0, END)
        # f = open(filename, 'r')
        # textPad.insert(1.0, f.read())

        with gzip.open(filename, 'rb') as input:
            with io.TextIOWrapper(input, encoding='utf-8') as dec:
                content2 = dec.read()
                textPad.insert(1.0, content2)
        filename.close()


def addtime():
    now = time.strftime('%Y-%m-%d %H:%M:%S')
    textPad.insert(1.0, now)


def new():
    global root, filename, textPad
    root.title('untitle')
    filename = None
    textPad.delete(1.0, END)


def save():
    global filename
    try:
        f = open(filename, 'w')
        msg = textPad.get(1.0, 'end')
        f.write(msg)
        f.close()
    except:
        saveas()


def save2():
    global filename
    try:
        # f = open(filename, 'w')
        msg = textPad.get(1.0, 'end')
        outfilename = asksaveasfilename(initialfile='untitle.note', defaultextension='.note')
        with gzip.open(outfilename, 'wb') as output:
            with io.TextIOWrapper(output, encoding='utf-8') as enc:
                enc.write(msg)
        os.system('file -b --mime {}'.format(outfilename))
        # f.write(msg)
        outfilename.close()
    except:
        # saveas()
        pass


def saveas():
    f = asksaveasfilename(initialfile='untitle.txt', defaultextension='.txt')
    global filename
    filename = f
    fh = open(f, 'w')
    msg = textPad.get(1.0, END)
    fh.write(msg)
    fh.close()
    root.title('Rogabet Notepad   ' + os.path.basename(f))


def cut():
    global textPad
    textPad.event_generate('<<Cut>>')


def copy():
    global textPad
    textPad.event_generate('<<Copy>>')


def paste():
    global textPad
    textPad.event_generate('<<Paste>>')


def undo():
    global textPad
    textPad.event_generate('<<Undo>>')


def redo():
    global textPad
    textPad.event_generate('<<Redo>>')


def select_all():
    global textPad
    textPad.tag_add('sel', '1.0', END)


# 'sel', '1.0', 'end'


def find():
    global root
    t = Toplevel(root)
    t.title('Find-Replace')
    # 设置窗口大小
    t.geometry('440x100+200+250')
    t.transient(root)
    v1 = StringVar()
    Label(t, text='Find:').grid(row=0, column=0, sticky='e')
    Label(t, text='Replace as:').grid(row=1, column=0)
    Entry(t, width=28, textvariable=v1).grid(row=1, column=1)

    v = StringVar()
    e = Entry(t, width=28, textvariable=v)  # 替换

    e.grid(row=0, column=1, padx=2, pady=2, sticky='we')
    e.focus_set()
    c = IntVar()

    Checkbutton(t, text='Match case', variable=c).grid(row=2, column=1, sticky='e')
    Button(t, text='find all', command=lambda: search(v.get(), c.get(), textPad, t, e)).grid(row=0, column=3,
                                                                                             sticky='e' + 'w', padx=2,
                                                                                             pady=2)
    Button(t, text='replace', command=lambda: mytihuan(v1.get(), v.get())).grid(row=1, column=3, padx=2, pady=2)

    # tihuantext = Text(t, width=3, height=2)

    def close_search():
        textPad.tag_remove('match', '1.0', END)
        t.destroy()

        t.protocol('WM_DELETE_WINDOW', close_search)


def mytihuan(tihuanwenben, yuanshiwenben):
    showinfo('hello', "replaced")
    find_data = yuanshiwenben.strip()
    replace_data = tihuanwenben.strip()
    data = textPad.get(1.0, END)
    print("finddata" + find_data)
    data = data.replace(find_data, replace_data)
    textPad.delete(1.0, END)
    textPad.insert(1.0, data)

    # textPad.mark_set(data)

    def search(needle, cssnstv, textPad, t, e):
        textPad.tag_remove('match', '1.0', END)
        count = 0
        if needle:
            pos = '1.0'
            while True:
                pos = textPad.search(needle, pos, nocase=cssnstv, stopindex=END)
                if not pos: break
                # lastpos=0
                lastpos = pos + str(len(needle))
                # print(str(len(needle))+"-----"+needle)
                textPad.tag_add('match', pos, lastpos)
                count += 1
                pos = lastpos
            textPad.tag_config('match', foreground='yellow', background='green')
            e.focus_set()
            t.title(str(count) + 'matched')

    def close_search():
        textPad.tag_remove('match', '1.0', END)
        t.destroy()
        t.protocol('WM_DELETE_WINDOW', close_search)


def search(needle, cssnstv, textPad, t, e):
    textPad.tag_remove('match', '1.0', END)
    count = 0
    if needle:
        pos = '1.0'
        while True:
            pos = textPad.search(needle, pos, nocase=cssnstv, stopindex=END)
            if not pos: break
            lastpos = pos + str(len(needle))
            textPad.tag_add('match', pos, lastpos)
            count += 1
            pos = lastpos
        textPad.tag_config('match', foreground='yellow', background='green')
        e.focus_set()
        t.title(str(count) + 'matched')


def popup(event):
    global editmenu
    editmenu.tk_popup(event.x_root, event.y_root)


def quit():
    exit()




def toCn():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'en'  # 源语言
    toLang = 'zh'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)

    '''
appid = '20151113000005349'
secretKey = 'osubCEzlGjzvw8qdQc41'
--
appid = '20191219000367626'
secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    '''


def getmacsn():
    macsncmd = 'cat |system_profiler SPHardwareDataType > macsn.txt'
    os.popen(macsncmd)

    time.sleep(1)

    with open("macsn.txt", "rt") as rdfile:
        getsn = rdfile.read()
    textPad.insert(1.0, getsn)


def toEn():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'

    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'zh'  # 源语言
    toLang = 'en'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)


def toJp():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'zh'  # 源语言
    toLang = 'jp'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)


def fromJp():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'jp'  # 源语言
    toLang = 'zh'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)


def toCht():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'zh'  # 源语言
    toLang = 'cht'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)


def fromCht():
    appid = '20191219000367626'
    secretKey = 'nRpvbOF3PgUaLXkaKrIi'
    httpClient = None
    myurl = '/api/trans/vip/translate'
    content = textPad.get(1.0, END).replace('\n', '')
    old = textPad.get(1.0, END)
    q = content
    fromLang = 'cht'  # 源语言
    toLang = 'zh'  # 翻译后的语言
    salt = random.randint(32768, 65536)
    sign = appid + q + str(salt) + secretKey
    sign = hashlib.md5(sign.encode()).hexdigest()
    myurl = myurl + '?appid=' + appid + '&q=' + urllib.parse.quote(
        q) + '&from=' + fromLang + '&to=' + toLang + '&salt=' + str(
        salt) + '&sign=' + sign
    httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
    httpClient.request('GET', myurl)
    # response是HTTPResponse对象
    response = httpClient.getresponse()
    jsonResponse = response.read().decode("utf-8")  # 获得返回的结果，结果为json格式
    js = json.loads(jsonResponse)  # 将json格式的结果转换字典结构
    dst = str(js["trans_result"][0]["dst"])  # 取得翻译后的文本结果
    print(dst)  # 打印结果
    out2 = old + '\n' + dst + '\n'
    textPad.delete(0.0, 'end')
    textPad.insert('end', out2)



def play2():
    import webbrowser
    play_url = 'https://wx.qq.com/'

    webbrowser.open(play_url)


def youdao():
    import webbrowser
    play_url = 'http://fanyi.youdao.com/'
    webbrowser.open(play_url)


def ciba():
    import webbrowser
    play_url = 'http://www.iciba.com/fy'
    webbrowser.open(play_url)


import binascii


def Unicode2HexStr():
    st = textPad.get(1.0, END) + 'rogabet'
    # 0072006f00670061006200650074000a
    Hex_Str = ""
    for i in range(0, len(st)):
        Hex_Str += (hex(ord(st[i])).replace('0x', '').zfill(4))

    textPad.delete(0.0, 'end')
    textPad.insert('end', Hex_Str)


def HexStr2Unicode():
    st2 = textPad.get(1.0, END)
    Unicde_Str = ""
    for i in range(0, len(st2) // 4):
        chr(int(st2[i * 4:i * 4 + 4], 16))
        Unicde_Str += chr(int(st2[i * 4:i * 4 + 4], 16))
    Unicde_Str = Unicde_Str.replace('rogabet', '')
    textPad.delete(0.0, 'end')
    textPad.insert('end', Unicde_Str)


import shutil, datetime


def tpzwz():
    import webbrowser
    play_url = 'https://ocr.wdku.net/'
    webbrowser.open(play_url)
def yyzwz2():
    import webbrowser
    play_url = 'https://app.xunjiepdf.com/voice2text'
    webbrowser.open(play_url)

def ynote():
    import webbrowser
    play_url = 'https://note.youdao.com/web'
    webbrowser.open(play_url)


def wkweb():
    import webbrowser
    play_url = 'https://baike.baidu.com/'
    webbrowser.open(play_url)


def wdzh():
    import webbrowser
    play_url = 'https://cn.office-converter.com/'
    webbrowser.open(play_url)


def qrcodes():
    import webbrowser
    play_url = 'https://cli.im/'
    webbrowser.open(play_url)


def fly():
    import webbrowser
    play_url = 'https://app.sharetome.com/login'
    webbrowser.open(play_url)




def fanyisuoming():
    readdme = '''
Author: Rogabet.Luo
Date:2021-6-11
Many functions rely on the network, please connect to the network.

	'''
    showinfo('info', readdme)


import gzip
import io


def jiamidabao():
    showinfo('info', 'dict.note in same folder')
    content = textPad.get(1.0, END)
    outfilename = 'dict.note'
    with gzip.open(outfilename, 'wb') as output:
        with io.TextIOWrapper(output, encoding='utf-8') as enc:
            enc.write(content)
    os.system('file -b --mime {}'.format(outfilename))


def jiemijiebao():
    showinfo('info', ' dict.note in same folder')
    with gzip.open('dict.note', 'rb') as input:
        with io.TextIOWrapper(input, encoding='utf-8') as dec:
            content2 = dec.read()
            textPad.insert(1.0, content2)


def familyChanged(event):
    f = Font(family=familyVar.get(),weight=weightVar.get(),size=sizeVar.get())
    textPad.configure(font=f)


def weightChanged(event):
    f = Font(weight=weightVar.get(),family=familyVar.get(),size=sizeVar.get())
    textPad.configure(font=f)


def sizeSelected(event):
    f = Font(size=sizeVar.get(),family=familyVar.get(),weight=weightVar.get())
    textPad.configure(font=f)

def convertUTF8ToANSI():
	global oldfile
	oldfile=askopenfilename(defaultextension='.txt')
	global onewfile
	newfile = asksaveasfilename(initialfile='-ansi-code.txt', defaultextension='.txt')
	fp_ansi = open(newfile, 'wb')  # 转码后输出的文件
	fp_utf8 = open(oldfile, 'rb')  # 待转码的文件
	data = ""
	data = fp_utf8.read()
	data = data.decode('utf-8')  # 以二进制格式读入的数据需要先转码（转为内部码）才可以继续转换
	data = data.encode('mbcs', errors='ignore')  # 关键，将内部码再次编码
	fp_ansi.write(data)  # 写入文件
	fp_ansi.close()
	fp_utf8.close()
from tkinter.colorchooser import *
def bgcolor():
	mycolor=askcolor()
	textPad.config(bg=mycolor[1])

import csv
def addcsv():
	global addcsvfile
	addcsvfile = askopenfilename(defaultextension='.csv')
	csv_file=open(addcsvfile,'r')
	with csv_file:
		read_csv=csv.reader(csv_file)
		for row in read_csv:
			textPad.insert('end', row)
			textPad.insert('end', '\n')
from pdfplumber import open as openpd
#import pdfplumber
def addpdf():
	global addpdffile
	addpdffile = askopenfilename(defaultextension='.pdf')
	with openpd(addpdffile) as pdf:
		numpages = len(pdf.pages)
		for i in range(0, numpages):
			read_page = pdf.pages[i]
			#print(read_page.extract_text())
			textPad.insert('end', read_page.extract_text())
			textPad.insert('end', '\n-----\n')
from openpyxl import load_workbook
#import openpyxl
def addxlsx():
	global addxlsxfile
	addxlsxfile = askopenfilename(defaultextension='.xlsx')
	wb = load_workbook(addxlsxfile)
	# 获取workbook中所有的表格
	sheets = wb.sheetnames
	print(sheets)
	# 循环遍历所有sheet
	for i in range(len(sheets)):
		sheet = wb[sheets[i]]

		print('\n\nNo' + str(i + 1) + 'sheet: ' + sheet.title + '->>>')
		textPad.insert('end', '\n\nNo' + str(i + 1) + 'sheet: ' + sheet.title + '->>>')
		for r in range(1, sheet.max_row + 1):
			if r == 1:
				partt01='\n' + ''.join(
					[str(sheet.cell(row=r, column=c).value).ljust(17) for c in range(1, sheet.max_column + 1)])

				print('\n' + ''.join(
					[str(sheet.cell(row=r, column=c).value).ljust(17) for c in range(1, sheet.max_column + 1)]))
				textPad.insert('end',partt01)
				textPad.insert('end', '\n')
			else:
				partt02 =''.join([str(sheet.cell(row=r, column=c).value).ljust(20) for c in range(1, sheet.max_column + 1)])

				print(
					''.join([str(sheet.cell(row=r, column=c).value).ljust(20) for c in range(1, sheet.max_column + 1)]))
				textPad.insert('end', partt02)
				textPad.insert('end', '\n')

def newpptx():
	# python-pptx #add it in pycharm to import pptx
	import pptx
	global pptxfile

	pptxfile = askopenfilename(defaultextension='.pptx')
	presentation = pptx.Presentation(pptxfile)
	results = []
	for slide in presentation.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				for paragraph in shape.text_frame.paragraphs:
					part = []
					for run in paragraph.runs:
						part.append(run.text)
					results.append(''.join(part))
	results = [line for line in results if line.strip()]
	textPad.insert('end', results)

import pyAesCrypt
from os import stat,remove

def aesone():
	bufferSize=64*1024
	password="rogabet"
	global aesonefile
	aesonefile = askopenfilename(defaultextension='.txt')
	with open(aesonefile,'rb')as fin:
		with open(aesonefile+'.aes', 'wb')as fout:
			pyAesCrypt.encryptStream(fin,fout,password,bufferSize)
	encFileSize=stat(aesonefile+'.aes').st_size
def unaesone():
	bufferSize = 64 * 1024
	password = "rogabet"
	global unaesonefile
	unaesonefile = askopenfilename(defaultextension='.aes')
	encFileSize = stat(unaesonefile).st_size
	unaesonefileout = asksaveasfilename(initialfile='untitle.txt', defaultextension='.txt')
	with open(unaesonefile, 'rb')as fin:
		with open(unaesonefileout, 'wb')as fout:
			try:
				pyAesCrypt.decryptStream(fin, fout, password, bufferSize,encFileSize)
			except ValueError:
				remove(aesonefile)

import json
from threading import Timer
from wxpy import *
import requests
import urllib.parse



def superfind():
	aa=entryvar.get()
	zz=entryvar2.get()
	print(aa)
	print(zz)
	super= re.findall(r"%s(.+?)%s"%(aa,zz), textPad.get(1.0, END))

	print('super=', super)  #
	for i in super:
		print(i)
		super2 = str(i)

		print('super2=', super2)
		textPad.tag_config(super2,foreground='red')  #sn4= abbcc', 'avvff

		def search(text_widget, keyword, tag):
			pos = '1.0'
			while True:
				idx = text_widget.search(keyword, pos, END)
				if not idx:
					break
				pos = '{}+{}c'.format(idx, len(keyword))
				text_widget.tag_add(tag, idx, pos)

		search(textPad, super2,super2)
import docx
#pip install python-docx
def adddocx():
	global adddocxfile
	adddocxfile = askopenfilename(defaultextension='.docx')
	wdoc=docx.Document(adddocxfile)
	docxnumparas = len(wdoc.paragraphs)
	for i in range(0, docxnumparas):
		readdocx_para =wdoc.paragraphs[i].text
		textPad.insert('end', readdocx_para)
		textPad.insert('end', '\n')
def send_mail2():
	neirongmail = textPad.get(1.0, END)

	fromaddr = frommvar3.get()
	mailpass2 =mailpassvar.get()
	mailsmtp=mailsmtpvar.get()
	mailtitle=mailtitlevar.get()
	towho=getmvar3.get()
	import smtplib
	from email.mime.text import MIMEText
	from email.header import Header
	# from_addr = '328994875@qq.com'
	# password = 'yfbnmgmbjbc' # 发信方的信息：发信邮箱，QQ 邮箱授权码
	# to_addr = '664287094@qq.com' # 收信方邮箱
	from_addr = fromaddr
	password =mailpass2
	smtp_server = mailsmtp
	# 邮箱正文内容，第一个参数为内容，第二个参数为格式(plain 为纯文本)，第三个参数为编码
	msg = MIMEText(neirongmail, 'plain', 'utf-8')
	# 邮件头信息
	msg['From'] = Header(from_addr)
	msg['To'] = Header(towho)
	msg['Subject'] = Header(mailtitle)
	## 开启发信服务，这里使用的是加密传输
	server = smtplib.SMTP_SSL(smtp_server)
	server.connect(smtp_server, 465)
	# 登录发信邮箱
	server.login(from_addr, password)
	# 发送邮件
	server.sendmail(from_addr, towho, msg.as_string())
	# 关闭服务器
	server.quit()

	showinfo('hello', "Mail sent out")


import json
from threading import Timer
from wxpy import *
import requests
import urllib.parse
def send_wxnews():

	wxname = nickwxvar.get()

	if wxname=='':
		showinfo('hello', "Please Fill the nickname / remark of your friend in the blank space on the right of 'wechat msg to 'button,and then click'wechat msg to', and scan the QR code with your wechat in cell phone for authorization , and the text of the editing area will be sent out")
	else:
		global wxcontent
		wxcontent=textPad.get(1.0, END)
		bot = Bot()  # 连接微信,会出现一个登陆微信的二维码
		my_friend = bot.friends().search(wxname)[0]
		my_friend.send(wxcontent)
from tkinter.ttk import Button,OptionMenu,Frame,Label,Combobox
root = Tk()
root.title('Rogabet Notepad for UOS')
root.geometry('1024x555+100+100')

toolbar = Frame(root, relief=RAISED, borderwidth=1)
toolbar.pack(side=TOP, fill=X, padx=2, pady=1)

familyVar=StringVar()
familyFamily=("Arial","Times","Courier","Arial","Song")
familyVar.set(familyFamily[3])
family=ttk.OptionMenu(toolbar,familyVar,*familyFamily,command=familyChanged)
family.pack(side=LEFT,pady=2)

weightVar=StringVar()
weightFamily=("normal","normal","bold")
weightVar.set(weightFamily[0])
weight=OptionMenu(toolbar,weightVar,*weightFamily,command=weightChanged)
weight.pack(side=LEFT,pady=2)

sizeVar=IntVar()
size= ttk.Combobox(toolbar,textvariable=sizeVar,width=4)
sizeFamily=[x for x in range(10,73)]
size["values"]=sizeFamily
zitidx=size.current(6)
size.bind("<<ComboboxSelected>>",sizeSelected)
size.pack(side=LEFT)



# A Z search start
#https://my.oschina.net/u/2245781/blog/661533
glabel = Label(toolbar, text='A**Z:')
glabel.pack(side=LEFT,pady=2)
entryvar = StringVar()
gentry = Entry(toolbar, textvariable=entryvar,width=26)
gentry.pack(side=LEFT,pady=2)
entryvar2 = StringVar()
gentry = Entry(toolbar, textvariable=entryvar2,width=26)
gentry.pack(side=LEFT,pady=2)
gbutton = Button(toolbar, command=superfind, text='Find',width=5)
gbutton.pack(side=LEFT,pady=2)


#wechat send msg start
sendwxbutton = Button(toolbar, command=send_wxnews, text='Wechat msg to:',width=13)
sendwxbutton.pack(side=LEFT,pady=2)
nickwxvar = StringVar()
gentry = Entry(toolbar, textvariable=nickwxvar,width=16)
gentry.pack(side=LEFT,pady=2)


#forsendmail
toolbar2 = ttk.Frame(root, relief=RAISED, borderwidth=1)
toolbar2.pack(side=TOP, fill=X, padx=2, pady=1)
glabel2 = Label(toolbar2, text='From:')
glabel2.pack(side=LEFT, pady=2)
frommvar3 = StringVar()
gentry = Entry(toolbar2, textvariable=frommvar3, width=22)
gentry.pack(side=LEFT, pady=2)
glabel3 = Label(toolbar2, text='password:')
glabel3.pack(side=LEFT, pady=2)
mailpassvar = StringVar()
gentry = Entry(toolbar2, textvariable=mailpassvar, width=12)
gentry.pack(side=LEFT, pady=2)
glabel4 = Label(toolbar2, text='smtp:')
glabel4.pack(side=LEFT, pady=2)
mailsmtpvar = StringVar()
gentry = Entry(toolbar2, textvariable=mailsmtpvar, width=14)
gentry.pack(side=LEFT, pady=2)
glabel5 = Label(toolbar2, text='title:')
glabel5.pack(side=LEFT, pady=2)
mailtitlevar = StringVar()
gentry = Entry(toolbar2, textvariable=mailtitlevar, width=14)
gentry.pack(side=LEFT, pady=2)
glabel6 = Label(toolbar2, text='To:')
glabel6.pack(side=LEFT, pady=2)
getmvar3 = StringVar()
gentry = Entry(toolbar2, textvariable=getmvar3, width=22)
gentry.pack(side=LEFT, pady=2)
gbutton3 = Button(toolbar2, command=send_mail2, text='SendMail',width=9)
gbutton3.pack(side=LEFT, pady=2)


menubar = Menu(root)
filemenu = Menu(menubar)
filemenu.add_command(label='0 New file', command=new)
filemenu.add_command(label='1 Open txt', command=myopen)
filemenu.add_command(label='2 Save as txt', command=save)
filemenu.add_command(label='3 Open note', command=myopen2)
filemenu.add_command(label='4 Save as note', command=save2)
filemenu.add_command(label='5 Save as', command=saveas)
filemenu.add_command(label='6 UTF8->ANSI', command=convertUTF8ToANSI)
filemenu.add_command(label = 'Create qrcode', command=qrout,underline=0)
filemenu.add_command(label='Quit', command=root.quit)
menubar.add_cascade(label='File ', menu=filemenu)

editmenu = Menu(menubar)
editmenu.add_command(label='Undo', accelerator='Ctrl+Z', command=undo)
editmenu.add_command(label='Redo', accelerator='Ctrl+Y', command=redo)
editmenu.add_command(label='Cut', accelerator='Ctrl+X', command=cut)
editmenu.add_command(label='Copy', accelerator='Ctrl+C', command=copy)
editmenu.add_command(label='Paste', accelerator='Ctrl+V', command=paste)
editmenu.add_command(label='Select all', accelerator='Ctrl+A', command=select_all)
editmenu.add_command(label='Background color', command=bgcolor)
editmenu.add_separator()
editmenu.add_command(label='0 Find or replace', command=find)
editmenu.add_command(label='1 Insert date and time', command=addtime)
#editmenu.add_command(label='2 Insert this pc hardinfo', command=getmacsn)
editmenu.add_command(label='3 Insert text of pdf file', command=addpdf)
editmenu.add_command(label ='4 Insert text of pptx file', command=newpptx)
editmenu.add_command(label ='4 Insert text of docx file', command=adddocx)
editmenu.add_command(label='5 Insert text of xlsx (Excel) file', command=addxlsx)

menubar.add_cascade(label='Edit ', menu=editmenu)

toolmenu = Menu(menubar)
toolmenu.add_command(label='0 Ciba Online', command=ciba)
toolmenu.add_command(label='1 Youdao Online', command=youdao)
toolmenu.add_command(label='2 Cn to En ', command=toEn)
toolmenu.add_command(label='3 En to Cn ', command=toCn)
toolmenu.add_command(label='4 Chs -> Cht ', command=toCht)
toolmenu.add_command(label='5 Cht -> Chs ', command=fromCht)
toolmenu.add_command(label='6 Cn to Jpn', command=toJp)
toolmenu.add_command(label='7 Jpn to Cn ', command=fromJp)
menubar.add_cascade(label='Translate ', menu=toolmenu)
root.config(menu=menubar)

playmenu = Menu(menubar)
playmenu.add_command(label='1 Get text from Picture', command=tpzwz)
playmenu.add_command(label='2 Create qrcode online', command=qrcodes)
playmenu.add_command(label='3 Convert files tools', command=wdzh)
playmenu.add_command(label='4 Audio to words online',  command=yyzwz2)
playmenu.add_command(label='5 Youdao Cloud Note', command=ynote)
playmenu.add_command(label='6 Baidu Baike', command=wkweb)
playmenu.add_command(label='7 Wechat Online', command=play2)
menubar.add_cascade(label='Online-Tools', menu=playmenu)

codemenu = Menu(menubar)
codemenu.add_command(label='1 Encrypt ', command=Unicode2HexStr)
codemenu.add_command(label='2 Decrypt', command=HexStr2Unicode)
codemenu.add_command(label='3 AES Encrypt',  command=aesone)
codemenu.add_command(label='4 AES Decrypt',  command=unaesone)
menubar.add_cascade(label='Encrypt', menu=codemenu)

about = Menu(menubar)
about.add_command(label='Version', command=fanyisuoming)
menubar.add_cascade(label='About', menu=about)


textPad = Text(root, width=90, height=90, selectforeground="black", undo=True, font=90)
textPad.pack(expand=YES, fill=BOTH)
scroll = Scrollbar(textPad, width=5)
textPad.config(yscrollcommand=scroll.set)
scroll.config(command=textPad.yview)
scroll.pack(side=RIGHT, fill=Y)

textPad.bind('<Control-N>', new)
textPad.bind('<Control-n>', new)
textPad.bind('<Control-O>', myopen)
textPad.bind('<Control-o>', myopen)
textPad.bind('<Control-S>', save)
textPad.bind('<Control-s>', save)
textPad.bind('<Control-A>', select_all)
textPad.bind('<Control-a>', select_all)
textPad.bind('<Control-F>', find)
textPad.bind('<Control-f>', find)
textPad.bind('<Button-3>', popup)
root.mainloop()
